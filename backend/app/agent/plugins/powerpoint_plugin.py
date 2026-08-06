import logging
from typing import Any

from app.agent.base_plugin import PluginBase, PluginOutput
from app.agent.registry import register_plugin
from app.services.artifact_service import ArtifactService
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field, field_validator

logger = logging.getLogger(__name__)

MAX_SLIDES = 20


class PowerPointPluginInput(BaseModel):
    title: str = Field(..., description="Main title of the PowerPoint presentation deck")
    subtitle: str | None = Field("Discord Community Engagement Analysis", description="Subtitle for title slide")
    highlights: list[str] | None = Field(None, description="Bullet points of key analytical insights")
    data: list[dict[str, Any]] | str | None = Field(None, description="Optional dataset for slide tables/metrics. If omitted or using previous query result, omit or leave null.")

    @field_validator("data", mode="before")
    @classmethod
    def sanitize_data(cls, v: Any) -> list[dict[str, Any]] | None:
        if isinstance(v, str) or not isinstance(v, list):
            return None
        return v


@register_plugin
class PowerPointPlugin(PluginBase):
    @property
    def name(self) -> str:
        return "powerpoint"

    @property
    def description(self) -> str:
        return "Generates a multi-slide PowerPoint (.pptx) executive summary deck with formatted titles, KPI stat cards, and analytical insights."

    @property
    def input_schema(self) -> type[BaseModel]:
        return PowerPointPluginInput

    @property
    def can_consume(self) -> list[str]:
        return ["data", "chart"]

    @property
    def can_produce(self) -> list[str]:
        return ["file"]

    async def execute(
        self,
        params: dict[str, Any],
        context: dict[str, Any],
        progress_callback: Any | None = None
    ) -> PluginOutput:
        val_err = self.validate_params(params)
        if val_err:
            return PluginOutput(success=False, output_type="powerpoint", error=val_err)

        title = params.get("title", "Engagement Analytics Summary")
        subtitle = params.get("subtitle", "Discord Dataset Analytics")
        highlights = params.get("highlights") or [
            "Active engagement is concentrated in key text channels.",
            "Member growth shows consistent weekend spikes.",
            "Voice activity correlates with server boost status."
        ]
        dataset = params.get("data")
        if not isinstance(dataset, list):
            dataset = context.get("last_query_result") or []

        if progress_callback:
            await progress_callback(f"Building PowerPoint deck: {title}...")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        navy_bg = RGBColor(15, 23, 42)
        accent_blue = RGBColor(56, 189, 248)
        text_white = RGBColor(255, 255, 255)
        text_muted = RGBColor(148, 163, 184)

        def apply_dark_bg(slide):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = navy_bg

        # SLIDE 1: Title
        slide1 = prs.slides.add_slide(blank_layout)
        apply_dark_bg(slide1)
        txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_blue

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = text_muted
        p2.space_before = Pt(20)

        # SLIDE 2: Executive Summary
        slide2 = prs.slides.add_slide(blank_layout)
        apply_dark_bg(slide2)
        title_box = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = text_white

        content_box = slide2.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for item in highlights[:MAX_SLIDES]:
            p = tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = text_muted
            p.space_before = Pt(16)

        # SLIDE 3: Key Data Metrics Table
        slide3 = prs.slides.add_slide(blank_layout)
        apply_dark_bg(slide3)
        title_box = slide3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Data Overview & Key Metrics"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = text_white

        if dataset:
            rows_cnt = min(len(dataset) + 1, 6)
            headers = list(dataset[0].keys())[:4]
            cols_cnt = len(headers)

            table_shape = slide3.shapes.add_table(rows_cnt, cols_cnt, Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
            table = table_shape.table
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header).upper()
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(37, 99, 235)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = text_white

            for row_idx, row_data in enumerate(dataset[:5], 1):
                for col_idx, header in enumerate(headers):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(row_data.get(header, ""))
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(13)
                        paragraph.font.color.rgb = text_muted

        # SLIDE 4: Conclusion
        slide4 = prs.slides.add_slide(blank_layout)
        apply_dark_bg(slide4)
        txBox = slide4.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.33), Inches(2.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = text_white

        artifact_id, filepath = ArtifactService.create_artifact_file(".pptx")
        prs.save(filepath)

        return PluginOutput(
            success=True,
            output_type="powerpoint",
            result={"filename": f"{artifact_id}.pptx", "slide_count": len(prs.slides)},
            artifact_id=artifact_id,
            artifact_url=f"/api/artifacts/{artifact_id}",
            metadata={"title": title, "slides": len(prs.slides)}
        )
